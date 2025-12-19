from fastapi import FastAPI, UploadFile, Form
from fastapi.responses import FileResponse
from typing import List
from pptx import Presentation
from pdf2image import convert_from_path
import pandas as pd
import shutil
import os
import uuid

app = FastAPI()

SHEET_URL = "https://docs.google.com/spreadsheets/d/17XiGo7Q4VFJYAZ4ZydIwLOC1SxUXthZB5J8Jb_mYX-U/export?format=csv"

BASE_DIR = os.getcwd()
TEMPLATE_DIR = os.path.join(BASE_DIR, "template")
os.makedirs(TEMPLATE_DIR, exist_ok=True)

@app.get("/")
def root():
    return {"status": "Backend is running"}

# =========================
# งานที่ 1 : สร้าง PPT
# =========================
@app.post("/generate-ppt")
async def generate_ppt(
    name1: str = Form(""),
    name2: str = Form(""),
    name3: str = Form(""),
    name4: str = Form(""),
    name5: str = Form("")
):
    names = [name1, name2, name3, name4, name5]

    df = pd.read_csv(SHEET_URL)

    prs = Presentation("template/template.pptx")

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for i, name in enumerate(names, start=1):
                    shape.text = shape.text.replace(f"{{{{Name{i}}}}}", name)

    output_name = f"output_{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(BASE_DIR, output_name)
    prs.save(output_path)

    return FileResponse(
        output_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="result.pptx"
    )

# =========================
# งานที่ 2 : PDF → รูป
# =========================
@app.post("/pdf-to-image")
async def pdf_to_image(pdfs: List[UploadFile]):
    results = []

    for pdf in pdfs:
        pdf_path = os.path.join(BASE_DIR, pdf.filename)
        with open(pdf_path, "wb") as buffer:
            shutil.copyfileobj(pdf.file, buffer)

        images = convert_from_path(pdf_path)

        for i, img in enumerate(images, start=1):
            img_name = f"{pdf.filename}_p{i}.png"
            img_path = os.path.join(BASE_DIR, img_name)
            img.save(img_path)
            results.append(img_name)

    return {
        "message": "PDF converted",
        "images": results
    }
